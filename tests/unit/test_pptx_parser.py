import pytest
from pptx import Presentation

from acervo.core.exceptions import ArquivoParseError
from acervo.ingestion.parsers import pptx_parser


def test_extrai_texto_de_slide(tmp_path):
    caminho = tmp_path / "exemplo.pptx"
    apresentacao = Presentation()
    layout = apresentacao.slide_layouts[6]
    slide = apresentacao.slides.add_slide(layout)
    caixa = slide.shapes.add_textbox(0, 0, 100, 100)
    caixa.text_frame.text = "Regressão linear simples"
    apresentacao.save(str(caminho))

    blocos = pptx_parser.parse(caminho)

    assert len(blocos) == 1
    assert blocos[0].explicacao == "Regressão linear simples"


def test_pptx_corrompido_gera_erro_tratado(tmp_path):
    caminho = tmp_path / "corrompido.pptx"
    caminho.write_bytes(b"nao e um pptx valido")

    with pytest.raises(ArquivoParseError):
        pptx_parser.parse(caminho)
