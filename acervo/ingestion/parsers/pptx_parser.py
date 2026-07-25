"""Extrai texto de apresentações .pptx, um bloco por slide."""
from pathlib import Path

from pptx import Presentation
from pptx.exc import PackageNotFoundError

from acervo.core.exceptions import ArquivoParseError
from acervo.ingestion.parsers.base import BlocoBruto


def parse(caminho: Path) -> list[BlocoBruto]:
    try:
        apresentacao = Presentation(str(caminho))
    except (PackageNotFoundError, OSError, KeyError, ValueError) as e:
        raise ArquivoParseError(str(caminho), "pptx_parser", e) from e

    blocos: list[BlocoBruto] = []
    for ordem, slide in enumerate(apresentacao.slides):
        textos = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texto_shape = shape.text_frame.text.strip()
                if texto_shape:
                    textos.append(texto_shape)

        texto_slide = "\n".join(textos).strip()
        if not texto_slide:
            continue

        blocos.append(BlocoBruto(
            ordem=ordem,
            titulo=f"Slide {ordem + 1}",
            explicacao=texto_slide,
            codigo=None,
            linguagem=None,
        ))

    return blocos
